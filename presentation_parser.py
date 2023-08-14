from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
from typing import List


def extract_text_from_shape(shape) -> str:
    """
    Extracts the text from a shape, including its nested shapes.

    Args:
        shape: The shape object.

    Returns:
        str: The extracted text.
    """
    text = ""

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text += " ".join(run.text.strip() for run in paragraph.runs)
            text += "\n"

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            text += extract_text_from_shape(sub_shape)

    return text


def extract_text_from_slide(slide) -> str:
    """
    Extracts the text from all shapes in a slide, including nested shapes.

    Args:
        slide: The slide object.

    Returns:
        str: The extracted text from the slide.
    """
    slide_text = ""

    for shape in slide.shapes:
        slide_text += extract_text_from_shape(shape)

    return slide_text.strip()


def extract_text(path_to_presentation: str) -> List[str]:
    """
    Extracts the text from each shape in each slide of a PowerPoint presentation.

    Args:
        path_to_presentation (str): The file path to the PowerPoint presentation.

    Returns:
        List of strings, where each string represents the text from a slide.
    """
    try:
        prs = Presentation(path_to_presentation)
    except Exception as e:
        print(f"Error loading the presentation: {e}")
        return [""]
    presentation_text = []

    for slide in prs.slides:
        slide_text = extract_text_from_slide(slide)
        if slide_text:
            presentation_text.append(slide_text)

    return presentation_text
